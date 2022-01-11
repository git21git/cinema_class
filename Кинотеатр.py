import datetime
import random

from pptx import Presentation

COMMANDS = ('add_network', 'n', 'add_cinema', 'add_hall', 'add_movie', 'check_movie', 'show_hall',
            'help', 'exit', 'buy_ticket', 'c', 'h', 'm', 'cm', 'sh', 'bt', 'e',
            'create_prs', 'prs')

cinemas = {}
networks = {}


def interface():
    """Функция вывода подсказки"""
    with open('data/intro.txt', mode='r', encoding='utf8') as f:
        print(f.read())


def check_cinema(cinema):
    """Функция проверки корректности введенного кинотеатра"""
    if cinemas.get(cinema) is None:
        print(f'WARMING!! Кинотеатра с именем "{cinema}" не существует')
        print('Имеющиеся кинотеатры:')
        [print(cinemas[cinema]) for cinema in cinemas]
        return None
    return cinemas[cinema]


def check_network(nw):
    """Функция проверки корректности введенной сети(при добавлении кинотеатра)"""
    if networks.get(nw) is None:
        print(f'WARMING!! Сети Кинотеатров "{nw}" не существует.')
        print('Возможно, вы имели ввиду одну из этих сетей кинотеатров:')
        [print(networks[nq]) for nq in networks]
        return None
    return networks[nw]


def check_hall(cinema, hall_number):
    """Функция проверки корректности введенного номера зала в кинотеатре"""
    if cinema is None:
        return
    try:
        assert int(hall_number) > 0
        return int(hall_number) - 1
    except Exception:
        print(f'WARMING!! В кинотеатре "{cinema.name}" сети {cinema.network.name} такой зал отсутствует')
        return None


def check_movie(name, number):
    """Функция для проверки существования фильма в прокате и подходящих сеансов"""
    if number == 0:
        num = 1
    else:
        num = number
    free_movies = []
    is_looking_movie = False

    for cinema in cinemas:
        for i, hall in enumerate(cinemas[cinema].halls, 1):
            for j, movie in enumerate(hall.movies, 1):
                if movie.name != name:
                    continue
                is_looking_movie = True

                if check_is_free(movie.halls_plan, num):
                    free_movies.append((cinemas[cinema], hall, movie))

    if not is_looking_movie:
        print('WARMING!! Такого фильма нет в прокате')
        return

    if not free_movies:
        print('WARMING!! Нет подходящих сеансов')
    else:
        print('Найдены подходящие сеансы:')
        for i, (cinema, hall, movie) in enumerate(free_movies, 1):
            print(f'{i}. Кинотеатр "{cinema.name}", зал номер {cinema.halls.index(hall) + 1}')
            print(f'{len(str(i)) * " "}  фильм {movie.name}, сеанс {movie.start}-{movie.end}')


def check_is_free(hall_plan, num):
    """Функция для проверки существования фильма в прокате и подходящих сеансов"""
    for row in hall_plan:
        places = ''.join([el if el == ' .' else '+-' for el in row])
        if '+-' * num in places:
            return True
    return False


class Cinema:
    """Класс Кинотеатра. Содержит информацию о своем названии, списке залов
        Имеет методы: Добавление залов, Вывод зала, Печать количества залов"""

    def __init__(self, name, network):
        self.network = network
        self.name = name
        self.halls = []

    def __str__(self):
        return f'Кинотеатр "{self.name}": {len(self.halls)} залов'

    def __getitem__(self, item):
        return self.halls[item]

    def add_hall(self, hall):
        self.halls.append(hall)


class Cinema_Network():
    def __init__(self, name):
        self.name = name
        self.cinemas = []

    def __str__(self):
        return f'Сеть Кинотеатров "{self.name}": {len(self.cinemas)} кинотеатров'

    def __getitem__(self, item):
        return self.cinemas[item]

    def add_hall(self, cinema):
        self.cinemas.append(cinema)


class Hall:
    """Класс Кино-зала. Содержит информацию о названии КиноТеатра, плане зала, номере в КиноТеатре, Фильмах
            Имеет методы: Добавление плана, Вывод плана, Печать количества залов"""

    def __init__(self, cinema):
        self.cinema = cinema
        self.plan = self.make_halls_plan()
        self.num = len(self.cinema.halls) + 1
        self.movies = []

    def make_halls_plan(self):
        """Функция создает план зала в кинотеатре"""
        print('Введите размеры зала в 1 строчку через пробел: 0 < n < 16, 0 < m < 31',
              'Зал - прямоугольник M*N, где M - кол-во рядов,а N - кол-во мест в ряду.', sep='\n')

        n = m = 0
        while not (0 < n <= 15 and 0 < m <= 30):
            try:
                n, m = map(int, input('Введите размер зала: ').split())
            except ValueError:
                print('ERROR! Введены некорректные размеры')
                continue
            if not (0 < n <= 15 and 0 < m <= 30):
                print('ERROR! Введены недопустимые размеры')
        else:
            return [[str(i).rjust(2) for i in range(1, m + 1)] for _ in range(n)]

    def choose_movie(self, movies):
        """Функция выбора сеанса, при нескольких """
        if len(movies) == 1:
            return movies[0]
        print(f'В данном зале проходит несколько фильмов с названием "{movies[0].name}"')
        print(*[f'{i}. {film.name} {film.start}-{film.end}' for i, film in enumerate(movies, 1)],
              sep='\n')
        while True:
            try:
                user_ans = input('Выберите номер сеанса, на который хотите посмотреть места\n')
                assert 0 < int(user_ans) <= len(movies)
                return movies[int(user_ans) - 1]
            except Exception:
                print('ERROR! Введите корректный номер сеанса')

    def add_movie(self, movie):
        """Ф-ция для владельца: добавляем фильм в прокат"""
        for film in self.movies:
            if film == movie:
                print(f'WARMING!! В это время уже идет фильм "{film.name}" с {film.start} до {film.end}')
                return 'Error'
        self.movies.append(movie)
        print(f'Добавили фильм "{movie.name}" с {movie.start} до {movie.end}')

    def get_movie(self, movie):
        """Ф-ция для проверки наличия фильма"""
        free_movies = []
        for film in self.movies:
            if film.name == movie:
                free_movies.append(film)

        return self.choose_movie(free_movies) if free_movies else None


class Movie:
    """Класс фильма в прокате"""

    def __init__(self, hall):
        self.name = input('Введите название фильма\n')

        self.hall = hall
        self.halls_plan = hall.plan[:]
        self.num = len(self.hall.movies) + 1

        self.start = self.end = None
        self.set_time()

        self.orders = []

    def __eq__(self, other):
        """Функция для проверки пересечения времени двух фильмов в прокате"""
        return (other.start <= self.start < other.end or other.start < self.end < other.end or
                (self.start < other.start and self.end > other.end) or
                (self.start > other.start and self.end < other.end))

    def set_time(self):
        """Функция установки времени начала и конца фильма"""
        movie_time = []

        for i, time in enumerate((self.start, self.end)):
            while time is None:
                try:
                    if i == 0:
                        word = "начала"
                    else:
                        word = "конца"
                    time = datetime.time(*map(int, input(f'Введите время {word} фильма в формате "чч мм"\n').split()))
                    movie_time.append(time)
                except Exception:
                    print('ERROR! Введена некорректное время фильма. Формат ввода "чч мм".')

        self.start, self.end = sorted(movie_time)

    def show_hall(self):
        """Функция вывода зала(при покупке билетов)"""
        print(' ' * 3 + '-' * (len(self.halls_plan[0]) - 5) * 4)
        for i, row in enumerate(self.halls_plan, 1):
            print(str(i).ljust(2) + '-' + ' '.join(row))

    def buy_tickets(self):
        """Функция покупки билетов"""
        number = '6'
        while number not in '12345':
            number = input('Выберите количество билетов (от 1 до 5)\n')
        number = int(number)

        tickets = []
        for _ in range(number):
            print('Выберите места ("*" - забронированные места)')
            print('Укажите ряд и номер места через пробел (например 1 2)')
            self.show_hall()

            ticket = None
            while ticket is None:
                try:
                    n, m = map(int, input().split())
                    assert n > 0 and m > 0

                    if self.halls_plan[n - 1][m - 1] == ' *':
                        print('WARMING!! Данное место уже занято')

                    self.halls_plan[n - 1][m - 1] = ' *'
                    ticket = (n, m)
                    tickets.append(ticket)
                    print(f'Бронируем место {n} {m}')
                except (ValueError, AssertionError, IndexError):
                    print('ERROR! Выберите корректное место')

        self.show_hall()

        name = input('Введите ваше имя\n')
        self.orders.append((name, tickets, len(tickets)))
        print(f'Заказ № {len(self.orders)} принят!')


def generate_pres(cinema, movie):
    prs = Presentation()
    # создаем новый слайд со схемой для добавления изображений
    cinema = check_cinema(cinema)
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    slide.shapes.title.text = f"В прокат вышел фильм '{movie}'! \n " \
                              f"Смотрите только в Кинотеатре '{str(cinema.name)}'('{str(cinema.network.name)}')."
    subtitle = slide.placeholders[1]
    subtitle.text = "Тестовый текст"
    # добавляем изображение
    placeholder = slide.placeholders[1]
    num = random.choice(range(1, 6))
    placeholder.insert_picture(f'data/{num}.png')
    # сохраняем презентацию
    name = f'{movie}_prs.pptx'
    prs.save(name)
    return name


def main():
    """Основной цикл программы"""
    comd = input('Введите команду\n')
    while not comd or comd.split()[:1][0] not in COMMANDS:
        print('ERROR! Введена некорректная команда')
        comd = input('Введите команду\n')
    comd = comd.split()

    if comd[0] in ('add_cinema', 'c') and len(comd) == 3:
        nw = check_network(comd[1])
        name = comd[2]
        if nw is None:
            return True
        elif cinemas.get(name):
            print(f'WARMING!! Кинотеатр с именем "{name}" уже существует')
            return True
        print(f'В сеть "{comd[1]}" добавлен кинотеатр "{name}"')
        cinemas[name] = Cinema(name, nw)

    elif comd[0] in ('add_network', 'n') and len(comd) == 2:
        name = comd[1]
        if cinemas.get(name):
            print(f'WARMING!! Сеть кинотеатров "{name}" уже существует(')
            return True
        print(f'Добавлена сеть Кинотеатров "{name}"')
        networks[name] = Cinema_Network(name)

    elif comd[0] in ('add_hall', 'h') and len(comd) == 2:
        cinema = check_cinema(comd[1])
        if cinema is None:
            return True
        hall = Hall(cinema)
        print(f'Добавили зал в кинотеатр "{cinema.name}"')
        cinema.add_hall(hall)

    elif comd[0] in ('add_movie', 'm') and len(comd) == 3:
        cinema = check_cinema(comd[1])
        hall = check_hall(cinema, comd[2])
        if cinema is None or hall is None:
            return True
        movie = Movie(cinema[hall])
        if cinema[hall].add_movie(movie) == 'Error':
            return True

    elif comd[0] in ('check_movie', 'cm') and len(comd) == 2:
        print('У нас можно найти сеансы, на которых будет несколько мест рядом, если это Вам не нужно, введите 0')
        print('введите количество мест рядом (от 1 до 5)')
        number = input()
        while number not in '012345':
            print('ERROR! Введите число мест')
            number = input()
        check_movie(comd[1], int(number))

    elif comd[0] in ('show_hall', 'sh') and len(comd) == 4:
        cinema = check_cinema(comd[1])
        hall = check_hall(cinema, comd[2])
        if cinema is None or hall is None:
            return True
        movie = cinema[hall].get_movie(comd[3])
        if movie is None:
            print(f'WARMING!! Фильма с названием "{comd[3]}" нет в прокате в данном зале')
            return True
        print(f'Выводим места зала номер {hall + 1} в кинотеатре "{cinema}" на фильм "{movie.name}"')
        movie.show_hall()

    elif comd[0] in ('buy_ticket', 'bt') and len(comd) == 4:
        cinema = check_cinema(comd[1])
        hall = check_hall(cinema, comd[2])
        if cinema is None or hall is None:
            return True
        movie = cinema[hall].get_movie(comd[3])
        if movie is None:
            print('WARMING!! Данного фильма нет в прокате в данном зале')
            return True
        movie.buy_tickets()
    elif comd[0] in ('create_prs', 'prs') and len(comd) == 1:
        cinema, movie = input('Введите через пробел Кинотеатр и Фильм  ').split()
        if check_cinema(cinema):
            req = generate_pres(cinema, movie)

            print(f'Рекламный буклет "{req}" создан')

    elif comd[0] == 'exit' and len(comd) == 1:
        return False
    elif comd[0] == 'help' and len(comd) == 1:
        interface()
    else:
        print('ERROR! Введена некорректная команда')
    return True


if __name__ == '__main__':
    """Если наш знакомый откроет несколько сетей кинотеатров, 
        ему будет полезна билетная система с возможностью работы в нескольких сетях))"""
    print('Вас приветствует Билетная система "Кино".')
    interface()
    while main():
        pass
    print('Спасибо за использование нашей Билетной системы "Кино"')
